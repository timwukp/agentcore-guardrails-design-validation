"""Slide renderer for the AgentCore guardrails design-document decks.

Layout engine only — no content lives here. Two things it is careful about:

* **Fitting.** Every text block and table is measured before it is written and the
  font size is chosen so the block cannot run off the slide. Measurement counts a
  CJK codepoint as two half-em units and a Latin one as one, so the same slide
  spec fits in both languages at different sizes rather than overflowing in one.
* **East-Asian typefaces.** python-pptx only writes ``<a:latin>``; a run with
  Chinese text in it needs ``<a:ea>`` too or PowerPoint substitutes per-glyph.
  ``_style_run`` writes latin/ea/cs together.
"""

from __future__ import annotations

import math
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# theme
# --------------------------------------------------------------------------- #

NAVY = RGBColor(0x1B, 0x26, 0x31)
NAVY_L = RGBColor(0x2C, 0x3E, 0x50)
INK = RGBColor(0x1F, 0x24, 0x29)
MUTED = RGBColor(0x5F, 0x6B, 0x7A)
FAINT = RGBColor(0x8A, 0x96, 0xA3)
ORANGE = RGBColor(0xFF, 0x99, 0x00)
BLUE = RGBColor(0x14, 0x6E, 0xB4)
BLUE_L = RGBColor(0xE8, 0xF1, 0xFA)
GREEN = RGBColor(0x1E, 0x7E, 0x3C)
GREEN_L = RGBColor(0xE7, 0xF4, 0xEA)
RED = RGBColor(0xB3, 0x1D, 0x1D)
RED_L = RGBColor(0xFC, 0xEB, 0xEA)
AMBER = RGBColor(0x9A, 0x55, 0x00)
AMBER_L = RGBColor(0xFD, 0xF3, 0xE2)
LINE = RGBColor(0xD3, 0xDA, 0xE1)
BAND = RGBColor(0xF3, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

VERDICT_COLOR = {
    "TRUE": GREEN,
    "FALSE": RED,
    "INCONCLUSIVE": AMBER,
    "RECORDED": BLUE,
    "UNTESTABLE": MUTED,
    "NOT_TESTABLE": MUTED,
    "NOT EVALUABLE": MUTED,
}

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.62
TITLE_TOP = 0.40
BODY_TOP = 1.34
BODY_BOTTOM = 6.92
FOOT_TOP = 6.98
CONTENT_W = SLIDE_W - 2 * MARGIN

FONTS = {
    "en": {"latin": "Arial", "ea": "Arial", "mono": "Courier New"},
    "zh": {"latin": "Arial", "ea": "PingFang TC", "mono": "Courier New"},
}

# --------------------------------------------------------------------------- #
# inline markup
# --------------------------------------------------------------------------- #

_TOKEN = re.compile(r"(\*\*.+?\*\*|`[^`]+`)")
_VERDICT = re.compile(
    r"(?<![A-Za-z])(NOT_TESTABLE|NOT EVALUABLE|INCONCLUSIVE|UNTESTABLE|RECORDED|TRUE|FALSE)(?![A-Za-z])"
)


def _split_verdicts(text):
    """Colour sealed verdict words so they read as markers, not as prose."""
    out = []
    pos = 0
    for m in _VERDICT.finditer(text):
        if m.start() > pos:
            out.append((text[pos : m.start()], None))
        out.append((m.group(0), VERDICT_COLOR[m.group(0)]))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], None))
    return out


def inline_runs(text):
    """``**bold**`` / ``` `code` ``` → run specs. Verdict words get a colour."""
    runs = []
    for chunk in _TOKEN.split(text):
        if not chunk:
            continue
        if chunk.startswith("**") and chunk.endswith("**") and len(chunk) > 4:
            # A code span nested inside bold — `**Pin `boto3` ≥ 1.43.32**` — is the
            # common case in this document, so the inner text has to be tokenized
            # again; emitting it verbatim leaves literal backticks on the slide.
            for inner in _TOKEN.split(chunk[2:-2]):
                if not inner:
                    continue
                if inner.startswith("`") and inner.endswith("`") and len(inner) > 2:
                    runs.append({"t": inner[1:-1], "bold": True, "code": True, "color": None})
                    continue
                for t, c in _split_verdicts(inner):
                    runs.append({"t": t, "bold": True, "code": False, "color": c})
        elif chunk.startswith("`") and chunk.endswith("`") and len(chunk) > 2:
            runs.append({"t": chunk[1:-1], "bold": False, "code": True, "color": None})
        else:
            for t, c in _split_verdicts(chunk):
                runs.append({"t": t, "bold": False, "code": False, "color": c})
    return runs or [{"t": "", "bold": False, "code": False, "color": None}]


def vwidth(text):
    """Visual width in half-em units: CJK and full-width punctuation count double."""
    n = 0
    for ch in text:
        o = ord(ch)
        n += 2 if (o > 0x2E80 and not 0xFF61 <= o <= 0xFF9F) else 1
    return n


def plain(text):
    return re.sub(r"\*\*|`", "", text)


def wrapped_lines(text, width_in, size_pt, bold_factor=1.0):
    """How many rendered lines ``text`` needs in a box ``width_in`` wide."""
    cap = 144.0 * width_in / (size_pt * bold_factor)
    if cap <= 0:
        return 1
    total = 0
    for seg in plain(text).split("\n"):
        total += max(1, math.ceil(vwidth(seg) / cap))
    return total


# --------------------------------------------------------------------------- #
# low-level XML helpers
# --------------------------------------------------------------------------- #


def _sub(parent, tag, after_tags=()):
    el = parent.find(qn(tag))
    if el is not None:
        return el
    el = parent.makeelement(qn(tag), {})
    anchor = None
    # `after_tags` is ordered nearest-first: the FIRST one present is the correct
    # anchor. Taking the last match instead puts `a:cs` ahead of `a:ea`, which is
    # out of schema order (latin → ea → cs) and makes PowerPoint repair the file.
    for t in after_tags:
        found = parent.find(qn(t))
        if found is not None:
            anchor = found
            break
    if anchor is not None:
        anchor.addnext(el)
    else:
        parent.append(el)
    return el


def _dash(shape_line, val="dash"):
    ln = shape_line._get_or_add_ln()
    el = _sub(ln, "a:prstDash")
    el.set("val", val)


def _arrowhead(shape_line, kind="triangle", at="tail"):
    ln = shape_line._get_or_add_ln()
    el = _sub(ln, f"a:{at}End")
    el.set("type", kind)
    el.set("w", "med")
    el.set("len", "med")


# --------------------------------------------------------------------------- #
# deck
# --------------------------------------------------------------------------- #


class Deck:
    def __init__(self, lang, running_title=""):
        self.lang = lang
        self.f = FONTS[lang]
        self.running_title = running_title
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self._blank = self.prs.slide_layouts[6]
        self.section = ""
        self.n = 0
        # Layout is computed, not eyeballed, so anything that could not be made to fit
        # has to say so — a slide that quietly runs off the bottom looks deliberate.
        self.warnings = []
        self.last_title = ""

    # -- primitives -------------------------------------------------------- #

    def _style_run(self, run, size, color, bold=False, italic=False, mono=False):
        font = run.font
        font.size = Pt(size)
        font.bold = bold
        font.italic = italic
        font.color.rgb = color
        name = self.f["mono"] if mono else self.f["latin"]
        font.name = name
        rPr = run._r.get_or_add_rPr()
        ea = _sub(rPr, "a:ea", ("a:latin",))
        ea.set("typeface", self.f["mono"] if mono else self.f["ea"])
        cs = _sub(rPr, "a:cs", ("a:ea", "a:latin"))
        cs.set("typeface", name)

    def _para(
        self,
        tf,
        text,
        size,
        color=INK,
        bold=False,
        italic=False,
        align=PP_ALIGN.LEFT,
        space_before=0,
        space_after=4,
        line_spacing=1.06,
        indent=0,
        bullet=None,
        first=False,
    ):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space_before)
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if indent:
            p.paragraph_format.left_indent = Inches(indent) if hasattr(p, "paragraph_format") else None
            p._pPr.set("marL", str(int(Inches(indent))))
            p._pPr.set("indent", str(int(Inches(-0.19))))
        if bullet:
            runs = [{"t": bullet + " ", "bold": False, "code": False, "color": FAINT}]
            runs += inline_runs(text)
        else:
            runs = inline_runs(text)
        for spec in runs:
            # A raw newline inside <a:t> is not a line break in DrawingML — PowerPoint
            # collapses it. Each "\n" has to become its own <a:br/> between runs.
            for i, piece in enumerate(spec["t"].split("\n")):
                if i:
                    p._p.append(p._p.makeelement(qn("a:br"), {}))
                if not piece:
                    continue
                r = p.add_run()
                r.text = piece
                self._style_run(
                    r,
                    size,
                    spec["color"] or color,
                    bold=bold or spec["bold"],
                    italic=italic,
                    mono=spec["code"],
                )
        return p

    def _box(self, sl, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
        sp = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sp.adjustments[0] = radius
        tf = sp.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.07)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        return sp

    def _tb(self, sl, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        return tf

    # -- chrome ------------------------------------------------------------ #

    def _new(self, chrome=True, title=None, kicker=None):
        sl = self.prs.slides.add_slide(self._blank)
        if chrome:
            self.n += 1
            self.last_title = plain(title or "")
            if title is not None:
                tf = self._tb(sl, MARGIN, TITLE_TOP, CONTENT_W - 1.1, 0.82)
                if kicker:
                    self._para(tf, kicker, 11, BLUE, bold=True, space_after=2, first=True)
                    self._para(tf, title, 23, NAVY, bold=True, space_after=0, line_spacing=1.0)
                else:
                    self._para(tf, title, 24, NAVY, bold=True, space_after=0, first=True, line_spacing=1.0)
                rule = self._box(sl, MARGIN, BODY_TOP - 0.15, CONTENT_W, 0.022, fill=ORANGE)
                rule.text_frame.text = ""
            ftf = self._tb(sl, MARGIN, FOOT_TOP, CONTENT_W - 0.9, 0.3)
            self._para(ftf, self.running_title + ("  ·  " + self.section if self.section else ""), 8.5, FAINT, first=True)
            ntf = self._tb(sl, SLIDE_W - MARGIN - 0.85, FOOT_TOP, 0.85, 0.3)
            self._para(ntf, str(self.n), 8.5, FAINT, align=PP_ALIGN.RIGHT, first=True)
        return sl

    # -- slide kinds ------------------------------------------------------- #

    def title_slide(self, title, subtitle, meta_rows, badge=None):
        sl = self._new(chrome=False)
        self._box(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
        self._box(sl, 0, 0, 0.28, SLIDE_H, fill=ORANGE)
        tf = self._tb(sl, 1.25, 1.55, 10.6, 2.6)
        if badge:
            self._para(tf, badge, 12, ORANGE, bold=True, space_after=10, first=True)
        self._para(tf, title, 40, WHITE, bold=True, space_after=12, line_spacing=1.04, first=not badge)
        self._para(tf, subtitle, 17, RGBColor(0xB9, 0xC6, 0xD2), space_after=0, line_spacing=1.16)
        self._box(sl, 1.25, 4.55, 4.2, 0.02, fill=RGBColor(0x44, 0x55, 0x66))
        mtf = self._tb(sl, 1.25, 4.85, 10.6, 1.9)
        for i, (k, v) in enumerate(meta_rows):
            p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
            p.space_after = Pt(4)
            p.line_spacing = 1.05
            r = p.add_run()
            r.text = k + "   "
            self._style_run(r, 10.5, ORANGE, bold=True)
            for spec in inline_runs(v):
                r = p.add_run()
                r.text = spec["t"]
                self._style_run(r, 10.5, RGBColor(0xC8, 0xD3, 0xDD), bold=spec["bold"], mono=spec["code"])
        return sl

    def divider(self, number, title, blurb=None, items=()):
        sl = self._new(chrome=False)
        self.n += 1
        self._box(sl, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
        self._box(sl, 0, 0, 0.28, SLIDE_H, fill=ORANGE)
        tf = self._tb(sl, 1.25, 2.5, 10.4, 2.4)
        self._para(tf, number, 13, ORANGE, bold=True, space_after=8, first=True)
        self._para(tf, title, 33, WHITE, bold=True, space_after=14, line_spacing=1.05)
        if blurb:
            self._para(tf, blurb, 14, RGBColor(0xB9, 0xC6, 0xD2), space_after=0, line_spacing=1.2)
        if items:
            itf = self._tb(sl, 1.25, 5.25, 10.4, 1.4)
            for i, it in enumerate(items):
                self._para(itf, it, 11.5, RGBColor(0x9F, 0xAF, 0xBD), bullet="—", space_after=4, first=(i == 0))
        ntf = self._tb(sl, SLIDE_W - MARGIN - 0.85, FOOT_TOP, 0.85, 0.3)
        self._para(ntf, str(self.n), 8.5, RGBColor(0x55, 0x66, 0x77), align=PP_ALIGN.RIGHT, first=True)
        return sl

    def bullets(self, title, items, kicker=None, lead=None, note=None, columns=1, sizes=(14, 13, 12, 11, 10, 9.5, 9)):
        """``items``: list of ``(level, text)`` or ``(level, text, style)``.

        ``style`` ∈ {None, 'head', 'warn', 'good', 'muted'}.
        """
        sl = self._new(title=title, kicker=kicker)
        y = BODY_TOP
        avail_h = BODY_BOTTOM - BODY_TOP
        if lead:
            lh = self._flow_text(sl, lead, MARGIN, y, CONTENT_W, 12, MUTED, italic=False)
            y += lh + 0.14
        note_h = 0.0
        if note:
            note_h = wrapped_lines(note, CONTENT_W, 9) * 0.135 + 0.12
        avail_h = BODY_BOTTOM - y - note_h
        col_w = (CONTENT_W - 0.5 * (columns - 1)) / columns
        size = self._fit_bullets(items, col_w, avail_h, sizes, columns)
        self._write_bullets(sl, items, MARGIN, y, col_w, avail_h, size, columns)
        if note:
            self._note(sl, note, BODY_BOTTOM - note_h + 0.06)
        return sl

    def _flow_text(self, sl, text, x, y, w, size, color, italic=False, bold=False):
        h = wrapped_lines(text, w, size) * (size * 1.14 / 72.0) + 0.06
        tf = self._tb(sl, x, y, w, h)
        self._para(tf, text, size, color, italic=italic, bold=bold, space_after=0, first=True, line_spacing=1.14)
        return h

    def _bullet_metrics(self, items, size):
        out = []
        for it in items:
            lvl, text = it[0], it[1]
            style = it[2] if len(it) > 2 else None
            s = size * (1.0 if lvl == 0 else 0.93)
            if style == "head":
                s = size * 1.0
            out.append((lvl, text, style, s))
        return out

    def _fit_bullets(self, items, col_w, avail_h, sizes, columns):
        for size in sizes:
            total = 0.0
            per_col = [0.0] * columns
            ci = 0
            for lvl, text, style, s in self._bullet_metrics(items, size):
                w = col_w - (0.26 if lvl == 0 else 0.52)
                lines = wrapped_lines(text, w, s, bold_factor=1.04 if style == "head" else 1.0)
                h = lines * (s * 1.16 / 72.0) + (0.11 if lvl == 0 else 0.07)
                if style == "head":
                    h += 0.06
                total += h
            if columns == 1:
                if total <= avail_h:
                    return size
            else:
                if total / columns <= avail_h * 0.99:
                    return size
        return sizes[-1]

    def _write_bullets(self, sl, items, x, y, col_w, avail_h, size, columns):
        metrics = self._bullet_metrics(items, size)
        heights = []
        for lvl, text, style, s in metrics:
            w = col_w - (0.26 if lvl == 0 else 0.52)
            lines = wrapped_lines(text, w, s, bold_factor=1.04 if style == "head" else 1.0)
            h = lines * (s * 1.16 / 72.0) + (0.11 if lvl == 0 else 0.07)
            if style == "head":
                h += 0.06
            heights.append(h)
        # greedy column balance, never splitting a head from its first child
        target = sum(heights) / columns
        cols = [[] for _ in range(columns)]
        ci, acc = 0, 0.0
        for i, (m, h) in enumerate(zip(metrics, heights)):
            if ci < columns - 1 and acc > 0 and acc + h / 2 > target and m[0] == 0:
                ci += 1
                acc = 0.0
            cols[ci].append((m, h))
            acc += h
        tallest = max((sum(h for _, h in col) for col in cols), default=0.0)
        if tallest > avail_h + 0.02:
            self.warnings.append(
                f"slide {self.n} “{self.last_title[:60]}”: bullets overflow by "
                f"{tallest - avail_h:.2f}in at {size}pt — cut text or split the slide")
        for ci, col in enumerate(cols):
            cx = x + ci * (col_w + 0.5)
            cy = y
            for (lvl, text, style, s), h in col:
                indent = 0.0 if lvl == 0 else 0.26
                bw = col_w - indent
                color = INK
                bullet = "▪"
                if style == "head":
                    color = NAVY
                    bullet = None
                elif style == "warn":
                    color = RED
                    bullet = "!"
                elif style == "good":
                    color = GREEN
                    bullet = "✓"
                elif style == "muted":
                    color = MUTED
                    bullet = "–"
                elif lvl > 0:
                    bullet = "–"
                    color = RGBColor(0x3A, 0x42, 0x4A)
                tf = self._tb(sl, cx + indent, cy, bw, h)
                self._para(
                    tf,
                    text,
                    s,
                    color,
                    bold=(style == "head"),
                    bullet=bullet,
                    space_after=0,
                    first=True,
                    line_spacing=1.16,
                )
                cy += h

    def _note(self, sl, note, y):
        self._box(sl, MARGIN, y - 0.05, 0.028, wrapped_lines(note, CONTENT_W - 0.16, 9) * 0.135 + 0.1, fill=LINE)
        tf = self._tb(sl, MARGIN + 0.12, y, CONTENT_W - 0.12, 0.6)
        self._para(tf, note, 9, MUTED, italic=True, space_after=0, first=True, line_spacing=1.14)

    def table(
        self,
        title,
        headers,
        rows,
        col_ratios=None,
        kicker=None,
        lead=None,
        note=None,
        emphasis_col=None,
        row_colors=None,
        sizes=(12, 11, 10.5, 10, 9.5, 9, 8.5, 8, 7.5, 7),
        align=None,
    ):
        sl = self._new(title=title, kicker=kicker)
        y = BODY_TOP
        if lead:
            y += self._flow_text(sl, lead, MARGIN, y, CONTENT_W, 11.5, MUTED) + 0.12
        note_h = wrapped_lines(note, CONTENT_W, 9) * 0.135 + 0.16 if note else 0.0
        avail_h = BODY_BOTTOM - y - note_h
        ncols = len(headers)
        ratios = col_ratios or [1] * ncols
        tot = float(sum(ratios))
        widths = [CONTENT_W * r / tot for r in ratios]

        size = sizes[-1]
        row_hs = None
        for cand in sizes:
            hs = []
            for ri, row in enumerate(([headers] + rows)):
                s = cand * (0.95 if ri else 1.0)
                lines = 1
                for ci, cell in enumerate(row):
                    lines = max(lines, wrapped_lines(str(cell), widths[ci] - 0.16, s, bold_factor=1.06 if ri == 0 else 1.0))
                hs.append(max(0.26, lines * (s * 1.2 / 72.0) + 0.11))
            if sum(hs) <= avail_h:
                size, row_hs = cand, hs
                break
        if row_hs is None:
            size = sizes[-1]
            row_hs = []
            for ri, row in enumerate(([headers] + rows)):
                s = size * (0.95 if ri else 1.0)
                lines = 1
                for ci, cell in enumerate(row):
                    lines = max(lines, wrapped_lines(str(cell), widths[ci] - 0.16, s))
                row_hs.append(max(0.24, lines * (s * 1.2 / 72.0) + 0.09))
            if sum(row_hs) > avail_h + 0.02:
                self.warnings.append(
                    f"slide {self.n} “{self.last_title[:60]}”: table overflows by "
                    f"{sum(row_hs) - avail_h:.2f}in at {size}pt — lower the abridge budget, "
                    f"drop columns, or split with keep_rows")

        cy = y
        for ri, row in enumerate([headers] + rows):
            h = row_hs[ri]
            cx = MARGIN
            if ri == 0:
                self._box(sl, MARGIN, cy, CONTENT_W, h, fill=NAVY)
            else:
                fill = None
                if row_colors and (ri - 1) < len(row_colors) and row_colors[ri - 1]:
                    fill = row_colors[ri - 1]
                elif ri % 2 == 0:
                    fill = BAND
                if fill:
                    self._box(sl, MARGIN, cy, CONTENT_W, h, fill=fill)
                self._box(sl, MARGIN, cy, CONTENT_W, 0.012, fill=LINE)
            for ci, cell in enumerate(row):
                w = widths[ci]
                s = size * (0.95 if ri else 1.0)
                al = PP_ALIGN.LEFT
                if align and align[ci] == "c":
                    al = PP_ALIGN.CENTER
                tf = self._tb(sl, cx + 0.08, cy + 0.055, w - 0.16, h - 0.09, anchor=MSO_ANCHOR.TOP)
                self._para(
                    tf,
                    str(cell),
                    s,
                    WHITE if ri == 0 else INK,
                    bold=(ri == 0) or (emphasis_col is not None and ci == emphasis_col),
                    align=al,
                    space_after=0,
                    first=True,
                    line_spacing=1.2,
                )
                cx += w
            cy += h
        self._box(sl, MARGIN, cy, CONTENT_W, 0.012, fill=LINE)
        if note:
            self._note(sl, note, BODY_BOTTOM - note_h + 0.1)
        return sl

    def kpi(self, title, cards, items=(), kicker=None, lead=None, note=None):
        sl = self._new(title=title, kicker=kicker)
        y = BODY_TOP
        if lead:
            y += self._flow_text(sl, lead, MARGIN, y, CONTENT_W, 12, MUTED) + 0.14
        n = len(cards)
        gap = 0.18
        cw = (CONTENT_W - gap * (n - 1)) / n
        ch = 1.5
        for i, c in enumerate(cards):
            x = MARGIN + i * (cw + gap)
            accent = c.get("color", BLUE)
            self._box(sl, x, y, cw, ch, fill=c.get("fill", BAND), line=LINE, line_w=0.75)
            self._box(sl, x, y, cw, 0.045, fill=accent)
            tf = self._tb(sl, x + 0.14, y + 0.24, cw - 0.28, ch - 0.3)
            self._para(tf, c["value"], 33, accent, bold=True, space_after=2, first=True, line_spacing=1.0)
            self._para(tf, c["label"], 11.5, NAVY, bold=True, space_after=2, line_spacing=1.1)
            if c.get("sub"):
                self._para(tf, c["sub"], 9.5, MUTED, space_after=0, line_spacing=1.12)
        y += ch + 0.26
        note_h = wrapped_lines(note, CONTENT_W, 9) * 0.135 + 0.16 if note else 0.0
        if items:
            avail = BODY_BOTTOM - y - note_h
            size = self._fit_bullets(items, CONTENT_W, avail, (13, 12, 11.5, 11, 10.5, 10, 9.5, 9), 1)
            self._write_bullets(sl, items, MARGIN, y, CONTENT_W, avail, size, 1)
        if note:
            self._note(sl, note, BODY_BOTTOM - note_h + 0.1)
        return sl

    def twocol(self, title, left, right, kicker=None, lead=None, note=None):
        """``left``/``right``: ``(heading, colour, [(level, text), ...])``."""
        sl = self._new(title=title, kicker=kicker)
        y = BODY_TOP
        if lead:
            y += self._flow_text(sl, lead, MARGIN, y, CONTENT_W, 12, MUTED) + 0.14
        note_h = wrapped_lines(note, CONTENT_W, 9) * 0.135 + 0.16 if note else 0.0
        avail = BODY_BOTTOM - y - note_h
        cw = (CONTENT_W - 0.42) / 2
        size = min(
            self._fit_bullets(left[2], cw, avail - 0.52, (13, 12, 11.5, 11, 10.5, 10, 9.5, 9, 8.5), 1),
            self._fit_bullets(right[2], cw, avail - 0.52, (13, 12, 11.5, 11, 10.5, 10, 9.5, 9, 8.5), 1),
        )
        for i, (head, color, items) in enumerate((left, right)):
            x = MARGIN + i * (cw + 0.42)
            self._box(sl, x, y, cw, 0.4, fill=color)
            tf = self._tb(sl, x + 0.14, y + 0.09, cw - 0.28, 0.3)
            self._para(tf, head, 12.5, WHITE, bold=True, space_after=0, first=True)
            self._write_bullets(sl, items, x, y + 0.55, cw, avail - 0.55, size, 1)
        if note:
            self._note(sl, note, BODY_BOTTOM - note_h + 0.1)
        return sl

    # -- diagrams ---------------------------------------------------------- #

    def diagram(self, title, nodes, edges=(), groups=(), kicker=None, lead=None, note=None, legend=None):
        """Native-shape diagram. Coordinates are inches inside the body area.

        ``nodes``: ``{id: {x, y, w, h, text, fill, line, fg, size, shape}}``
        ``edges``: ``[{a, b, label, dash, color, side, waypoint}]``
        ``groups``: ``[{x, y, w, h, label, fill, line}]``
        """
        sl = self._new(title=title, kicker=kicker)
        y0 = BODY_TOP
        if lead:
            y0 += self._flow_text(sl, lead, MARGIN, y0, CONTENT_W, 11.5, MUTED) + 0.1
        for g in groups:
            gx, gy, gw, gh = MARGIN + g["x"], y0 + g["y"], g["w"], g["h"]
            self._box(sl, gx, gy, gw, gh, fill=g.get("fill", BAND), line=g.get("line", LINE), line_w=0.9,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03)
            if g.get("label"):
                tf = self._tb(sl, gx + 0.12, gy + 0.07, gw - 0.24, 0.26)
                self._para(tf, g["label"], g.get("size", 10), g.get("fg", NAVY), bold=True,
                           align=g.get("align", PP_ALIGN.LEFT), space_after=0, first=True)
        placed = {}
        for nid, n in nodes.items():
            x, yy, w, h = MARGIN + n["x"], y0 + n["y"], n["w"], n["h"]
            placed[nid] = (x, yy, w, h)
            shape = n.get("shape", MSO_SHAPE.ROUNDED_RECTANGLE)
            sp = self._box(sl, x, yy, w, h, fill=n.get("fill", WHITE), line=n.get("line", LINE),
                           line_w=n.get("line_w", 1.0), shape=shape, radius=0.06 if shape == MSO_SHAPE.ROUNDED_RECTANGLE else None)
            tf = sp.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            size = n.get("size")
            if size is None:
                size = 10.5
                while size > 6.5 and wrapped_lines(n["text"], w - 0.16, size) * (size * 1.16 / 72.0) > h - 0.08:
                    size -= 0.5
            self._para(tf, n["text"], size, n.get("fg", INK), bold=n.get("bold", False),
                       align=PP_ALIGN.CENTER, space_after=0, first=True, line_spacing=1.1)
        for e in edges:
            self._edge(sl, placed, e)
        if legend:
            lx = MARGIN
            ly = BODY_BOTTOM - 0.24
            for label, color in legend:
                self._box(sl, lx, ly + 0.05, 0.16, 0.1, fill=color)
                tf = self._tb(sl, lx + 0.22, ly, 2.4, 0.22)
                self._para(tf, label, 8.5, MUTED, space_after=0, first=True)
                lx += 0.28 + 0.09 * vwidth(label) * 8.5 / 12.0
        if note:
            nh = wrapped_lines(note, CONTENT_W, 9) * 0.135 + 0.16
            self._note(sl, note, BODY_BOTTOM - nh + 0.12)
        return sl

    @staticmethod
    def _anchor(rect, side):
        x, y, w, h = rect
        return {
            "l": (x, y + h / 2),
            "r": (x + w, y + h / 2),
            "t": (x + w / 2, y),
            "b": (x + w / 2, y + h),
        }[side]

    def _edge(self, sl, placed, e):
        ra, rb = placed[e["a"]], placed[e["b"]]
        side = e.get("side")
        if not side:
            ax, ay, aw, ah = ra
            bx, by, bw, bh = rb
            if bx >= ax + aw - 0.05:
                side = "rl"
            elif ax >= bx + bw - 0.05:
                side = "lr"
            elif by >= ay + ah - 0.05:
                side = "bt"
            else:
                side = "tb"
        p1 = self._anchor(ra, side[0])
        p2 = self._anchor(rb, side[1])
        color = e.get("color", RGBColor(0x7C, 0x8A, 0x97))
        pts = [p1] + list(e.get("waypoint", [])) + [p2]
        for i in range(len(pts) - 1):
            a, b = pts[i], pts[i + 1]
            c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(a[0]), Inches(a[1]), Inches(b[0]), Inches(b[1]))
            c.line.color.rgb = color
            c.line.width = Pt(e.get("w", 1.25))
            if e.get("dash"):
                _dash(c.line, e.get("dash") if isinstance(e.get("dash"), str) else "dash")
            if i == len(pts) - 2:
                _arrowhead(c.line, e.get("head", "triangle"))
        if e.get("label"):
            mid = pts[len(pts) // 2] if len(pts) > 2 else ((p1[0] + p2[0]) / 2, (p1[1] + p2[1]) / 2)
            lw = max(0.5, 0.055 * vwidth(plain(e["label"])) + 0.12)
            lines = max(1, e["label"].count("\n") + 1)
            lh = 0.14 * lines + 0.06
            off = e.get("loff", (0, 0))
            bx, by = mid[0] - lw / 2 + off[0], mid[1] - lh / 2 + off[1]
            # A long label centred on a waypoint can hang off the slide — the EN wording
            # of a feedback edge is routinely wider than the ZH one, so this cannot be
            # left to the caller's coordinates. Slide back inside the margins; the label
            # stays on its own arrow because these are the long spanning edges.
            bx = min(max(bx, MARGIN), max(MARGIN, SLIDE_W - MARGIN - lw))
            by = min(max(by, BODY_TOP), max(BODY_TOP, BODY_BOTTOM - lh))
            self._box(sl, bx, by, lw, lh, fill=WHITE, line=None)
            tf = self._tb(sl, bx, by + 0.01, lw, lh)
            self._para(tf, e["label"], e.get("lsize", 8.5), e.get("lcolor", MUTED), align=PP_ALIGN.CENTER,
                       space_after=0, first=True, line_spacing=1.05)

    # -- output ------------------------------------------------------------ #

    def save(self, path):
        self.prs.save(path)
        return len(self.prs.slides)
